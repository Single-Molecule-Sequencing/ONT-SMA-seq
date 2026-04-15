# report.py
# Generates presentation-ready reports from a completed SMA database.
#
# Four stages:
#   1. extract_summary()  — query DB for aggregate metrics  (stdlib only)
#   2. generate_figures()  — publication-quality plots        (matplotlib/seaborn)
#   3. generate_narrative() — AI-powered slide text           (google-generativeai, optional)
#   4. build_pptx() / build_markdown() — assemble output     (python-pptx or stdlib)

import json
import os
import sqlite3
import statistics
from pathlib import Path

# ---------------------------------------------------------------------------
# Stage 1 — Data Extraction & Summary
# ---------------------------------------------------------------------------

def extract_summary(db_path):
	"""
	Queries a completed SMA database and returns a dict of aggregate metrics.

	Returns:
		dict with keys: exp, target, read_count, readlen_*, q_bc_*, q_ld_*,
		ed_*, end_reason_counts, model_tier, model_ver, trim, mods.

	Raises:
		FileNotFoundError: If db_path does not exist.
		RuntimeError: If the database has no reads.
	"""
	db = Path(db_path)
	if not db.exists():
		raise FileNotFoundError(f"Database not found: {db_path}")

	conn = sqlite3.connect(db)
	conn.row_factory = sqlite3.Row
	c = conn.cursor()

	# --- Experiment metadata ---
	c.execute("SELECT * FROM Exp LIMIT 1")
	exp_row = c.fetchone()
	if not exp_row:
		conn.close()
		raise RuntimeError("No experiment metadata in DB. Was mkdb run?")

	exp = {
		"exp_id": exp_row["exp_id"],
		"flow_cell_id": exp_row["flow_cell_id"],
		"sample_id": exp_row["sample_id"],
		"alias": exp_row["alias"],
		"exp_desc": exp_row["exp_desc"],
	}

	# --- Target ---
	c.execute("SELECT tgt_id, tgt_reflen FROM Target LIMIT 1")
	tgt_row = c.fetchone()
	target = {
		"tgt_id": tgt_row["tgt_id"] if tgt_row else "unknown",
		"tgt_reflen": tgt_row["tgt_reflen"] if tgt_row else 0,
	}

	# --- Read count ---
	c.execute("SELECT COUNT(*) AS n FROM Reads")
	read_count = c.fetchone()["n"]
	if read_count == 0:
		conn.close()
		raise RuntimeError("Database contains no reads. Run ingest first.")

	# --- Numeric distributions ---
	c.execute("SELECT readlen, q_bc, q_ld, ed FROM Reads")
	readlens, q_bcs, q_lds, eds = [], [], [], []
	for row in c:
		readlens.append(row["readlen"])
		q_bcs.append(row["q_bc"])
		q_lds.append(row["q_ld"])
		eds.append(row["ed"])

	def _stats(values):
		return {
			"min": min(values),
			"max": max(values),
			"mean": round(statistics.mean(values), 2),
			"median": round(statistics.median(values), 2),
			"stdev": round(statistics.stdev(values), 2) if len(values) > 1 else 0.0,
		}

	# --- End-reason breakdown ---
	c.execute("SELECT ER, COUNT(*) AS cnt FROM Reads GROUP BY ER ORDER BY cnt DESC")
	end_reason_counts = {row["ER"]: row["cnt"] for row in c}

	# --- Mod string ---
	c.execute(
		"SELECT DISTINCT r.mod_bitflag, m.mods "
		"FROM Reads r JOIN Mods m ON r.mod_bitflag = m.mod_bitflag"
	)
	mod_rows = c.fetchall()
	mods_list = [{"bitflag": r["mod_bitflag"], "label": r["mods"]} for r in mod_rows]

	# --- Model info (take most common) ---
	c.execute(
		"SELECT model_tier, model_ver, trim, COUNT(*) AS cnt "
		"FROM Reads GROUP BY model_tier, model_ver, trim ORDER BY cnt DESC LIMIT 1"
	)
	model_row = c.fetchone()

	conn.close()

	return {
		"exp": exp,
		"target": target,
		"read_count": read_count,
		"readlen": _stats(readlens),
		"q_bc": _stats(q_bcs),
		"q_ld": _stats(q_lds),
		"ed": _stats(eds),
		"end_reason_counts": end_reason_counts,
		"mods": mods_list,
		"model_tier": model_row["model_tier"] if model_row else "?",
		"model_ver": model_row["model_ver"] if model_row else "?",
		"trim": model_row["trim"] if model_row else "?",
	}


# ---------------------------------------------------------------------------
# Stage 2 — Figure Generation  (requires matplotlib, seaborn)
# ---------------------------------------------------------------------------

def _check_analysis_deps():
	"""Raises ImportError with install hint if matplotlib/seaborn are missing."""
	try:
		import matplotlib  # noqa: F401
		import seaborn  # noqa: F401
	except ImportError as exc:
		raise ImportError(
			"Figure generation requires matplotlib and seaborn. "
			"Install with:  pip install ont-sma-seq[analysis]"
		) from exc


def generate_figures(db_path, out_dir):
	"""
	Generates publication-quality PNG figures from a completed SMA database.

	Produces:
		- readlen_hist.png   — Read-length distribution
		- quality_dist.png   — q_bc and q_ld violin/box plots
		- ed_vs_readlen.png  — Edit distance vs. read length scatter
		- end_reason.png     — End-reason bar chart

	Args:
		db_path: Path to the SQLite database.
		out_dir: Directory to write PNGs into (created if needed).

	Returns:
		list[Path]: Paths to generated figures.
	"""
	_check_analysis_deps()
	import matplotlib
	matplotlib.use("Agg")
	import matplotlib.pyplot as plt
	import seaborn as sns

	out = Path(out_dir)
	out.mkdir(parents=True, exist_ok=True)

	conn = sqlite3.connect(db_path)
	c = conn.cursor()

	c.execute("SELECT readlen, q_bc, q_ld, ed, ER FROM Reads")
	rows = c.fetchall()
	conn.close()

	readlens = [r[0] for r in rows]
	q_bcs    = [r[1] for r in rows]
	q_lds    = [r[2] for r in rows]
	eds      = [r[3] for r in rows]
	ers      = [r[4] for r in rows]

	sns.set_theme(style="whitegrid", context="talk")
	figures = []

	# 1) Read-length histogram
	fig, ax = plt.subplots(figsize=(10, 6))
	ax.hist(readlens, bins=50, color="#4C72B0", edgecolor="white", linewidth=0.5)
	ax.set_xlabel("Read Length (bp)")
	ax.set_ylabel("Count")
	ax.set_title("Read-Length Distribution")
	p = out / "readlen_hist.png"
	fig.savefig(p, dpi=150, bbox_inches="tight")
	plt.close(fig)
	figures.append(p)

	# 2) Quality distributions (side-by-side violin plots)
	fig, axes = plt.subplots(1, 2, figsize=(12, 6))
	sns.violinplot(y=q_bcs, ax=axes[0], color="#55A868", inner="box")
	axes[0].set_title("Basecall Quality (q_bc)")
	axes[0].set_ylabel("Phred Q")
	sns.violinplot(y=q_lds, ax=axes[1], color="#C44E52", inner="box")
	axes[1].set_title("Levenshtein Quality (q_ld)")
	axes[1].set_ylabel("Phred Q")
	fig.suptitle("Quality Score Distributions", y=1.02)
	fig.tight_layout()
	p = out / "quality_dist.png"
	fig.savefig(p, dpi=150, bbox_inches="tight")
	plt.close(fig)
	figures.append(p)

	# 3) Edit distance vs. read length scatter
	fig, ax = plt.subplots(figsize=(10, 6))
	ax.scatter(readlens, eds, alpha=0.3, s=8, color="#8172B2")
	ax.set_xlabel("Read Length (bp)")
	ax.set_ylabel("Edit Distance")
	ax.set_title("Edit Distance vs. Read Length")
	p = out / "ed_vs_readlen.png"
	fig.savefig(p, dpi=150, bbox_inches="tight")
	plt.close(fig)
	figures.append(p)

	# 4) End-reason bar chart
	from collections import Counter
	er_counts = Counter(ers)
	labels = list(er_counts.keys())
	counts = list(er_counts.values())
	fig, ax = plt.subplots(figsize=(10, 6))
	bars = ax.barh(labels, counts, color="#CCB974", edgecolor="white")
	ax.set_xlabel("Count")
	ax.set_title("End-Reason Breakdown")
	ax.bar_label(bars, fmt="%d", padding=3)
	fig.tight_layout()
	p = out / "end_reason.png"
	fig.savefig(p, dpi=150, bbox_inches="tight")
	plt.close(fig)
	figures.append(p)

	return figures


# ---------------------------------------------------------------------------
# Stage 3 — Gemini-Powered Narrative  (optional)
# ---------------------------------------------------------------------------

_GEMINI_PROMPT = """\
You are a scientific report writer for Oxford Nanopore sequencing experiments.
Given the following JSON summary of a Single-Molecule-Accuracy (SMA-seq) run,
produce a concise presentation narrative in JSON format with these keys:

  "title_text":     A one-line title for the experiment slide deck.
  "overview_text":  A 2–3 sentence experiment overview.
  "quality_text":   A short paragraph interpreting the basecall (q_bc) and
                    Levenshtein (q_ld) quality metrics.
  "ed_text":        A short paragraph interpreting the edit-distance distribution.
  "end_reason_text": A short paragraph on the end-reason breakdown.
  "conclusion_text": A 2–3 sentence conclusion with key takeaways.
  "talking_points":  A list of 3–5 bullet strings for the presenter.

Return ONLY valid JSON — no markdown, no extra text.
"""


def _default_narrative(summary):
	"""Template-based narrative when Gemini is unavailable."""
	exp = summary["exp"]
	return {
		"title_text": f"SMA-seq Report — {exp['exp_id']}",
		"overview_text": (
			f"Experiment {exp['exp_id']} processed {summary['read_count']:,} reads "
			f"against target {summary['target']['tgt_id']} "
			f"({summary['target']['tgt_reflen']:,} bp). "
			f"Basecaller: tier {summary['model_tier']}, version {summary['model_ver']}."
		),
		"quality_text": (
			f"Basecall quality (q_bc): median {summary['q_bc']['median']}, "
			f"mean {summary['q_bc']['mean']} (σ {summary['q_bc']['stdev']}). "
			f"Levenshtein quality (q_ld): median {summary['q_ld']['median']}, "
			f"mean {summary['q_ld']['mean']} (σ {summary['q_ld']['stdev']})."
		),
		"ed_text": (
			f"Edit distance: median {summary['ed']['median']}, "
			f"mean {summary['ed']['mean']}, range [{summary['ed']['min']}–{summary['ed']['max']}]."
		),
		"end_reason_text": (
			"End-reason distribution: "
			+ ", ".join(
				f"{er} ({cnt:,})" for er, cnt in summary["end_reason_counts"].items()
			)
			+ "."
		),
		"conclusion_text": (
			f"This run produced {summary['read_count']:,} reads with a median "
			f"basecall quality of {summary['q_bc']['median']} and median "
			f"Levenshtein quality of {summary['q_ld']['median']}."
		),
		"talking_points": [
			f"Total reads: {summary['read_count']:,}",
			f"Target: {summary['target']['tgt_id']} ({summary['target']['tgt_reflen']:,} bp)",
			f"Median basecall quality: {summary['q_bc']['median']}",
			f"Median edit distance: {summary['ed']['median']}",
		],
	}


def generate_narrative(summary, gemini_api_key=None):
	"""
	Generates slide narrative text from a summary dict.

	If *gemini_api_key* is provided, calls the Gemini API for an AI-written
	narrative.  Otherwise falls back to a deterministic template.

	Args:
		summary: dict returned by extract_summary().
		gemini_api_key: Optional Google AI API key.

	Returns:
		dict with narrative text keyed by slide section.
	"""
	if not gemini_api_key:
		print("[report] No Gemini API key — using template narrative.")
		return _default_narrative(summary)

	try:
		from google import generativeai as genai
	except ImportError:
		print(
			"[report] google-generativeai not installed — using template narrative. "
			"Install with:  pip install ont-sma-seq[report]"
		)
		return _default_narrative(summary)

	print("[report] Generating AI narrative via Gemini...")
	genai.configure(api_key=gemini_api_key)
	model = genai.GenerativeModel("gemini-2.0-flash")

	# Send only aggregate stats, never raw sequences
	payload = json.dumps(summary, indent=2, default=str)
	prompt = _GEMINI_PROMPT + "\n\nSummary JSON:\n" + payload

	try:
		response = model.generate_content(prompt)
		text = response.text.strip()
		# Strip markdown fences if present
		if text.startswith("```"):
			text = text.split("\n", 1)[1]
		if text.endswith("```"):
			text = text.rsplit("```", 1)[0]
		narrative = json.loads(text)
		print("[report] Gemini narrative received.")
		return narrative
	except Exception as exc:
		print(f"[report] Gemini call failed ({exc}) — falling back to template.")
		return _default_narrative(summary)


# ---------------------------------------------------------------------------
# Stage 4a — PowerPoint Assembly  (requires python-pptx)
# ---------------------------------------------------------------------------

def build_pptx(summary, narrative, figures, out_path):
	"""
	Assembles a PowerPoint slide deck.

	Slides: Title → Overview → Quality → Edit Distance → End Reasons → Conclusion.

	Args:
		summary:   dict from extract_summary().
		narrative: dict from generate_narrative().
		figures:   list[Path] from generate_figures().
		out_path:  Output .pptx path.

	Returns:
		Path to the generated file.
	"""
	try:
		from pptx import Presentation
		from pptx.util import Inches, Pt
	except ImportError as exc:
		raise ImportError(
			"PPTX generation requires python-pptx. "
			"Install with:  pip install ont-sma-seq[report]"
		) from exc

	prs = Presentation()
	prs.slide_width = Inches(13.333)
	prs.slide_height = Inches(7.5)

	def _add_title_slide(title, subtitle):
		layout = prs.slide_layouts[0]  # Title Slide
		slide = prs.slides.add_slide(layout)
		slide.shapes.title.text = title
		slide.placeholders[1].text = subtitle

	def _add_content_slide(title, body_text, image_path=None):
		if image_path and Path(image_path).exists():
			# Use a blank layout and place elements manually
			layout = prs.slide_layouts[6]  # Blank
			slide = prs.slides.add_slide(layout)
			# Title box
			from pptx.util import Emu
			txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
			tf = txBox.text_frame
			p = tf.paragraphs[0]
			p.text = title
			p.font.size = Pt(28)
			p.font.bold = True
			# Image — centered
			slide.shapes.add_picture(
				str(image_path), Inches(1.5), Inches(1.2), Inches(10), Inches(5.0)
			)
			# Caption below image
			txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(6.4), Inches(12), Inches(0.9))
			tf2 = txBox2.text_frame
			tf2.word_wrap = True
			p2 = tf2.paragraphs[0]
			p2.text = body_text
			p2.font.size = Pt(14)
		else:
			layout = prs.slide_layouts[1]  # Title and Content
			slide = prs.slides.add_slide(layout)
			slide.shapes.title.text = title
			slide.placeholders[1].text = body_text

	# Build a lookup for figures by filename stem
	fig_map = {f.stem: f for f in figures} if figures else {}

	# Slide 1: Title
	_add_title_slide(
		narrative.get("title_text", f"SMA-seq Report — {summary['exp']['exp_id']}"),
		f"Experiment {summary['exp']['exp_id']}\n"
		f"Target: {summary['target']['tgt_id']} ({summary['target']['tgt_reflen']:,} bp)\n"
		f"Total Reads: {summary['read_count']:,}",
	)

	# Slide 2: Overview
	_add_content_slide(
		"Experiment Overview",
		narrative.get("overview_text", ""),
	)

	# Slide 3: Quality
	_add_content_slide(
		"Quality Metrics",
		narrative.get("quality_text", ""),
		fig_map.get("quality_dist"),
	)

	# Slide 4: Read Length
	_add_content_slide(
		"Read-Length Distribution",
		f"Median: {summary['readlen']['median']:,} bp  |  "
		f"Mean: {summary['readlen']['mean']:,} bp  |  "
		f"σ: {summary['readlen']['stdev']:,} bp",
		fig_map.get("readlen_hist"),
	)

	# Slide 5: Edit Distance
	_add_content_slide(
		"Edit Distance vs. Read Length",
		narrative.get("ed_text", ""),
		fig_map.get("ed_vs_readlen"),
	)

	# Slide 6: End Reasons
	_add_content_slide(
		"End-Reason Breakdown",
		narrative.get("end_reason_text", ""),
		fig_map.get("end_reason"),
	)

	# Slide 7: Conclusion & Talking Points
	points = narrative.get("talking_points", [])
	bullet_text = "\n".join(f"• {pt}" for pt in points) if points else ""
	_add_content_slide(
		"Conclusions",
		narrative.get("conclusion_text", "") + "\n\n" + bullet_text,
	)

	prs.save(str(out_path))
	return Path(out_path)


# ---------------------------------------------------------------------------
# Stage 4b — Markdown Report  (no extra deps)
# ---------------------------------------------------------------------------

def build_markdown(summary, narrative, figures, out_path):
	"""
	Assembles a Markdown report.

	Args:
		summary:   dict from extract_summary().
		narrative: dict from generate_narrative().
		figures:   list[Path] from generate_figures() (referenced as relative paths).
		out_path:  Output .md path.

	Returns:
		Path to the generated file.
	"""
	out = Path(out_path)
	fig_map = {f.stem: f for f in figures} if figures else {}

	lines = []
	lines.append(f"# {narrative.get('title_text', 'SMA-seq Report')}\n")

	lines.append("## Experiment Overview\n")
	lines.append(narrative.get("overview_text", "") + "\n")

	lines.append("| Metric | Value |")
	lines.append("|--------|-------|")
	lines.append(f"| Experiment ID | {summary['exp']['exp_id']} |")
	lines.append(f"| Target | {summary['target']['tgt_id']} ({summary['target']['tgt_reflen']:,} bp) |")
	lines.append(f"| Total Reads | {summary['read_count']:,} |")
	lines.append(f"| Model Tier | {summary['model_tier']} |")
	lines.append(f"| Model Version | {summary['model_ver']} |")
	lines.append(f"| Trim | {summary['trim']} |")
	lines.append("")

	lines.append("## Quality Metrics\n")
	lines.append(narrative.get("quality_text", "") + "\n")
	if "quality_dist" in fig_map:
		lines.append(f"![Quality Distributions]({fig_map['quality_dist']})\n")

	lines.append("## Read-Length Distribution\n")
	lines.append(
		f"Median: {summary['readlen']['median']:,} bp | "
		f"Mean: {summary['readlen']['mean']:,} bp | "
		f"σ: {summary['readlen']['stdev']:,} bp\n"
	)
	if "readlen_hist" in fig_map:
		lines.append(f"![Read-Length Distribution]({fig_map['readlen_hist']})\n")

	lines.append("## Edit Distance\n")
	lines.append(narrative.get("ed_text", "") + "\n")
	if "ed_vs_readlen" in fig_map:
		lines.append(f"![Edit Distance vs Read Length]({fig_map['ed_vs_readlen']})\n")

	lines.append("## End-Reason Breakdown\n")
	lines.append(narrative.get("end_reason_text", "") + "\n")
	if "end_reason" in fig_map:
		lines.append(f"![End-Reason Breakdown]({fig_map['end_reason']})\n")

	lines.append("## Conclusions\n")
	lines.append(narrative.get("conclusion_text", "") + "\n")
	points = narrative.get("talking_points", [])
	for pt in points:
		lines.append(f"- {pt}")
	lines.append("")

	out.write_text("\n".join(lines))
	return out


# ---------------------------------------------------------------------------
# Orchestrator
# ---------------------------------------------------------------------------

def run(db_path, out_dir=None, fmt="pptx", gemini_api_key=None):
	"""
	End-to-end report generation.

	Args:
		db_path:       Path to a completed SMA SQLite database.
		out_dir:       Output directory for report artifacts. Defaults to a
		               sibling ``report/`` directory next to the database.
		fmt:           Output format — ``"pptx"``, ``"md"``, or ``"html"``
		               (html is an alias for md).
		gemini_api_key: Optional Google AI API key for Gemini narrative.

	Returns:
		Path to the generated report file.
	"""
	db = Path(db_path)
	if out_dir is None:
		out_dir = db.parent / "report"
	out = Path(out_dir)
	out.mkdir(parents=True, exist_ok=True)
	fig_dir = out / "figures"

	print(f"[report] Database:  {db}")
	print(f"[report] Output:    {out}")
	print(f"[report] Format:    {fmt}")

	# Stage 1
	print("[report] Extracting summary...")
	summary = extract_summary(db_path)
	summary_path = out / "summary.json"
	summary_path.write_text(json.dumps(summary, indent=2, default=str))
	print(f"[report] Summary written to {summary_path}")

	# Stage 2
	figures = []
	try:
		print("[report] Generating figures...")
		figures = generate_figures(db_path, fig_dir)
		print(f"[report] {len(figures)} figures saved to {fig_dir}")
	except ImportError:
		print("[report] matplotlib/seaborn not available — skipping figures.")

	# Stage 3
	api_key = gemini_api_key or os.environ.get("GEMINI_API_KEY")
	narrative = generate_narrative(summary, gemini_api_key=api_key)

	# Stage 4
	fmt = fmt.lower().strip()
	if fmt == "pptx":
		report_file = out / f"report_{summary['exp']['exp_id']}.pptx"
		try:
			build_pptx(summary, narrative, figures, report_file)
		except ImportError:
			print("[report] python-pptx not available — falling back to Markdown.")
			fmt = "md"

	if fmt in ("md", "html"):
		report_file = out / f"report_{summary['exp']['exp_id']}.md"
		build_markdown(summary, narrative, figures, report_file)

	print(f"\n[report] Report generated: {report_file}")
	return report_file
